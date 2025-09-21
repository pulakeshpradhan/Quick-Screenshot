import os
import argparse
from pptx import Presentation
from pptx.util import Inches
from PIL import Image

# Set up argument parser
parser = argparse.ArgumentParser(description='Create a PowerPoint slideshow from images, scaled to fit the slide.')
parser.add_argument('--width', type=float, default=16, help='Width of the slides in inches.')
parser.add_argument('--height', type=float, default=9, help='Height of the slides in inches.')
args = parser.parse_args()

# Create a new presentation
prs = Presentation()
prs.slide_width = Inches(args.width)
prs.slide_height = Inches(args.height)
blank_slide_layout = prs.slide_layouts[6]

# Get the directory of the script
dir_path = os.path.dirname(os.path.realpath(__file__))

# List all files in the directory and filter for PNG images
image_files = [f for f in os.listdir(dir_path) if f.lower().endswith('.png')]

# ✅ Sort the images based on filename (YYYYMMDD_HHMMSS ensures correct order)
image_files.sort(key=lambda f: os.path.splitext(f)[0])

# Loop through each image and add it to a new slide
for image_name in image_files:
    slide = prs.slides.add_slide(blank_slide_layout)
    img_path = os.path.join(dir_path, image_name)

    # Get image dimensions using Pillow
    try:
        with Image.open(img_path) as img:
            img_width_px, img_height_px = img.size
    except Exception as e:
        print(f"Could not read image {image_name}: {e}")
        continue

    # Calculate the aspect ratios
    slide_aspect_ratio = prs.slide_width / prs.slide_height
    image_aspect_ratio = img_width_px / img_height_px

    # Scale the image to fit the slide while maintaining aspect ratio
    if image_aspect_ratio > slide_aspect_ratio:
        # Image is wider than the slide, so scale to slide width
        width = prs.slide_width
        height = int(width / image_aspect_ratio)
    else:
        # Image is taller than or same as the slide, so scale to slide height
        height = prs.slide_height
        width = int(height * image_aspect_ratio)

    # Center the image
    left = int((prs.slide_width - width) / 2)
    top = int((prs.slide_height - height) / 2)

    # Add the picture to the slide
    slide.shapes.add_picture(img_path, left, top, width=width, height=height)

# Save the presentation
output_filename = 'slide.pptx'
prs.save(output_filename)

print(f"Presentation '{output_filename}' created successfully with {len(image_files)} images!")
